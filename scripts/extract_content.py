"""
Extract text content from PDF or PPTX files.
Outputs JSON with extracted text per page/slide.
"""

import json
import sys
import os


def extract_from_pdf(file_path: str) -> list[dict]:
    import pdfplumber

    pages = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            tables = page.extract_tables() or []
            table_text = ""
            for table in tables:
                for row in table:
                    cells = [str(c) if c else "" for c in row]
                    table_text += " | ".join(cells) + "\n"
            combined = text
            if table_text:
                combined += "\n[表データ]\n" + table_text
            pages.append({
                "page": i + 1,
                "text": combined.strip(),
            })
    return pages


def extract_from_pptx(file_path: str) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        slides.append({
            "page": i + 1,
            "text": "\n".join(texts),
        })
    return slides


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_content.py <file_path>", file=sys.stderr)
        sys.exit(1)

    file_path = sys.argv[1]
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        result = extract_from_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        result = extract_from_pptx(file_path)
    else:
        print(f"Unsupported file type: {ext}", file=sys.stderr)
        sys.exit(1)

    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
