"""
Anime PPTX Generator
Creates a widescreen (16:9) PowerPoint presentation with anime images as full-bleed backgrounds.
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def create_pptx(data: dict) -> None:
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    blank_layout = prs.slide_layouts[6]  # blank layout

    for slide_info in data["slides"]:
        slide = prs.slides.add_slide(blank_layout)

        # Add image as full background
        img_path = slide_info["image_path"]
        slide.shapes.add_picture(
            img_path,
            Emu(0),
            Emu(0),
            slide_width,
            slide_height,
        )

        # Optionally add title text overlay
        if slide_info.get("show_title") and slide_info.get("topic"):
            # Semi-transparent dark bar at bottom
            left = Emu(0)
            top = slide_height - Inches(1.2)
            width = slide_width
            height = Inches(1.2)

            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                left,
                top,
                width,
                height,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
            shape.fill.fore_color.brightness = 0
            shape.line.fill.background()

            # Set transparency via XML
            from pptx.oxml.ns import qn
            solid_fill = shape.fill._fill
            a_solid_fill = solid_fill.find(qn("a:solidFill"))
            if a_solid_fill is not None:
                srgb = a_solid_fill.find(qn("a:srgbClr"))
                if srgb is not None:
                    alpha = srgb.makeelement(qn("a:alpha"), {})
                    alpha.set("val", "50000")  # 50% opacity
                    srgb.append(alpha)

            # Add title text
            txBox = slide.shapes.add_textbox(
                Inches(0.5),
                slide_height - Inches(1.1),
                slide_width - Inches(1),
                Inches(1.0),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_info["topic"]
            p.font.size = Pt(36)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT

    prs.save(data["output_path"])


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python create_pptx.py <json_path>", file=sys.stderr)
        sys.exit(1)

    json_path = sys.argv[1]
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    create_pptx(data)
    print(f"PPTX created: {data['output_path']}")
